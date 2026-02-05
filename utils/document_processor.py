import PyPDF2
from docx import Document
from pptx import Presentation

class DocumentProcessor:
    @staticmethod
    def extract_text(uploaded_file):
        if not uploaded_file: return ""
        name = uploaded_file.name.lower()
        
        try:
            if name.endswith('.pdf'):
                reader = PyPDF2.PdfReader(uploaded_file)
                return " ".join([p.extract_text() for p in reader.pages[:15]])
            elif name.endswith('.docx'):
                doc = Document(uploaded_file)
                return " ".join([p.text for p in doc.paragraphs])
            elif name.endswith('.pptx'):
                prs = Presentation(uploaded_file)
                return " ".join([shape.text for s in prs.slides for shape in s.shapes if hasattr(shape, "text")])
            elif name.endswith(('.png', '.jpg', '.jpeg')):
                return "[IMAGE_READY]"
        except Exception as e:
            return f"Processor Error: {e}"
        return "Unsupported Format"
